from pptx import Presentation
from pptx.util import Pt
from openai import OpenAI
import os
import tempfile
import time
from config import OPENAI_API_KEY

client = OpenAI(api_key=OPENAI_API_KEY)

def generate_slide_content(topic):
    prompt = f"""
    You are an academic assistant preparing a presentation.

    Topic: {topic}

    Provide:
    1. 3 concise bullet points explaining the topic clearly for students.
       (Make it simple and memorable like quick-glance notes.)

    2. A list of 4-5 important keywords.

    Format:
    Bullets:
    - point 1
    - point 2
    ...

    Keywords: keyword1, keyword2, ...
    """
    try:
        response = client.chat.completions.create(
            model="gpt-4",
            messages=[{"role": "user", "content": prompt}],
            temperature=0.5,
        )
        content = response.choices[0].message.content.strip()
        bullets, keywords = "", ""
        if "Bullets:" in content:
            parts = content.split("Bullets:")[1]
            if "Keywords:" in parts:
                bullets, keywords = parts.split("Keywords:")
        return bullets.strip(), keywords.strip()
    except Exception as e:
        return [f"Error: {e}"], ""

def generate_detailed_pptx(summary_text, descriptive=False, progress_callback=None):
    prs = Presentation()
    slide_layout = prs.slide_layouts[1]  # Title + Content

    lines = [line.strip() for line in summary_text.split("\n") if line.strip()]
    
    headings = {}
    current_heading = None

    for line in lines:
        if not line.startswith("-") and not line.startswith("*") and ":" not in line:
            current_heading = line
            headings[current_heading] = {"topics": [], "subtopics": []}
        elif ":" in line:
            headings[current_heading]["topics"].append(line.strip())
        else:
            headings[current_heading]["subtopics"].append(line.strip("-•* "))

    total_slides = sum([1 + len(v["topics"]) + (1 if v["subtopics"] else 0) for v in headings.values()])
    done = 0

    for heading, content in headings.items():
        # Slide 1: Heading title only
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = heading
        slide.placeholders[1].text = f"Overview of: {heading}"
        done += 1
        if progress_callback:
            progress_callback(done / total_slides)
        time.sleep(0.2)

        # Slide 2+: Topic slides with bullets + keywords
        for topic in content["topics"]:
            bullets, keywords = generate_slide_content(topic)
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = topic
            tf = slide.placeholders[1].text_frame
            tf.clear()

            bullet_lines = [line.strip("-•* ") for line in bullets.split("\n") if line.strip()]
            for b in bullet_lines:
                p = tf.add_paragraph()
                p.text = b
                p.level = 0
                p.font.size = Pt(16)

            if keywords:
                kw = tf.add_paragraph()
                kw.text = f"📌 Keywords: {keywords.strip()}"
                kw.level = 0
                kw.font.size = Pt(14)

            done += 1
            if progress_callback:
                progress_callback(done / total_slides)
            time.sleep(0.2)

        # Slide: All subtopics grouped (if any)
        if content["subtopics"]:
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.title.text = f"{heading} – Subtopics"
            tf = slide.placeholders[1].text_frame
            tf.clear()
            for st in content["subtopics"]:
                p = tf.add_paragraph()
                p.text = f"• {st}"
                p.level = 0
                p.font.size = Pt(14)
            done += 1
            if progress_callback:
                progress_callback(done / total_slides)
            time.sleep(0.2)

    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
        prs.save(tmp.name)
        return tmp.name

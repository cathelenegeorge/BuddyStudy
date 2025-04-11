from langchain.prompts import PromptTemplate
from langchain.chains import LLMChain
from langchain_openai import ChatOpenAI
from pptx import Presentation
import os
import tempfile
import re
import streamlit as st

llm = ChatOpenAI(model_name="gpt-4o", temperature=0.3)
flashcard_prompt = PromptTemplate(
    input_variables=["summary", "count"],
    template="""
You are a flashcard generator. From the summary below:

{summary}

Create exactly {count} well-formed flashcards using this format:

---
Q: <One-line question>
A: <One-line answer>
Explanation: <Brief explanation>

---

Important rules:
- DO NOT include any extra explanation outside the format.
- DO NOT add introduction or closing remarks.
- FOLLOW the format strictly.
"""
)

def generate_flashcards(summary, count, retries=2):
    chain = LLMChain(llm=llm, prompt=flashcard_prompt)

    for _ in range(retries):
        output = chain.run(summary=summary, count=count)
        pattern = r"Q:\s*(.*?)\s*A:\s*(.*?)\s*Explanation:\s*(.*?)(?=\n---|\Z)"
        flashcards = re.findall(pattern, output, re.DOTALL)

        if len(flashcards) >= int(count):
            return flashcards[:count]  # return only requested amount
    return flashcards  # return however many were valid

def export_flashcards_to_pptx(flashcards):
    prs = Presentation()
    for q, a, explanation in flashcards:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title, content = slide.shapes.title, slide.shapes.placeholders[1]
        title.text = f"Q: {q}"
        content.text = f"A: {a}\n\nExplanation: {explanation.strip()}"
    filepath = os.path.join(tempfile.gettempdir(), "flashcards.pptx")
    prs.save(filepath)
    return filepath
